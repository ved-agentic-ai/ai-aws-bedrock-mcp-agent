import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    BG_DARK = RGBColor(7, 9, 19)
    BG_CARD = RGBColor(20, 28, 48)
    TEXT_WHITE = RGBColor(248, 250, 252)
    TEXT_MUTED = RGBColor(148, 163, 184)
    ACCENT_AWS = RGBColor(255, 153, 0)
    ACCENT_BEDROCK = RGBColor(236, 72, 153)
    ACCENT_MCP = RGBColor(6, 182, 212)
    ACCENT_BLUE = RGBColor(59, 130, 246)

    def add_blank_slide():
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        # Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()
        return slide

    def add_header(slide, title_text, subtitle_text):
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.size = Pt(13)
            p2.font.color.rgb = ACCENT_MCP
            p2.space_before = Pt(4)

    # ==========================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================
    slide1 = add_blank_slide()
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "Enterprise Agentic AI Control Center"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT_AWS

    p2 = tf1.add_paragraph()
    p2.text = "Production-Grade Decoupled AI Architecture on AWS Bedrock & Model Context Protocol (MCP)"
    p2.font.size = Pt(20)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(14)

    p3 = tf1.add_paragraph()
    p3.text = "Amazon Bedrock (Nova Micro)  •  AWS Lambda  •  Model Context Protocol (MCP)  •  Amazon DynamoDB"
    p3.font.size = Pt(14)
    p3.font.color.rgb = ACCENT_MCP
    p3.space_before = Pt(24)

    # ==========================================
    # SLIDE 2: EXECUTIVE OVERVIEW
    # ==========================================
    slide2 = add_blank_slide()
    add_header(slide2, "Executive Overview & Architectural Shift", "Moving from monolithic prompt engineering to standardized Model Context Protocol (MCP) standards.")

    # Card 1: Traditional Monolithic AI
    card1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    card1.fill.solid()
    card1.fill.fore_color.rgb = BG_CARD
    card1.line.color.rgb = RGBColor(239, 68, 68)

    tf_c1 = card1.text_frame
    tf_c1.margin_left = Inches(0.3)
    tf_c1.margin_top = Inches(0.3)
    tf_c1.word_wrap = True
    
    p = tf_c1.paragraphs[0]
    p.text = "❌ Legacy Monolithic AI Architecture"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(239, 68, 68)

    points1 = [
        "• Hardcoded tool definitions tightly coupled to specific LLM providers.",
        "• High risk of vendor lock-in when migrating models.",
        "• Direct database access from prompts creates security vulnerabilities.",
        "• High idle infrastructure costs on dedicated virtual machines."
    ]
    for pt in points1:
        p_pt = tf_c1.add_paragraph()
        p_pt.text = pt
        p_pt.font.size = Pt(13)
        p_pt.font.color.rgb = TEXT_MUTED
        p_pt.space_before = Pt(12)

    # Card 2: Modern Decoupled Agentic AI
    card2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
    card2.fill.solid()
    card2.fill.fore_color.rgb = BG_CARD
    card2.line.color.rgb = ACCENT_MCP

    tf_c2 = card2.text_frame
    tf_c2.margin_left = Inches(0.3)
    tf_c2.margin_top = Inches(0.3)
    tf_c2.word_wrap = True

    p = tf_c2.paragraphs[0]
    p.text = "✅ Modern Agentic MCP Architecture"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_MCP

    points2 = [
        "• Standardized JSON-RPC protocol (tools/list & tools/call).",
        "• Zero Vendor Lock-In: Swap underlying LLMs in 1 config change.",
        "• Zero-Trust Isolation: MCP Lambda sandbox handles DB execution.",
        "• Pay-Per-Request Serverless: 100% idle cost optimization."
    ]
    for pt in points2:
        p_pt = tf_c2.add_paragraph()
        p_pt.text = pt
        p_pt.font.size = Pt(13)
        p_pt.font.color.rgb = TEXT_WHITE
        p_pt.space_before = Pt(12)

    # ==========================================
    # SLIDE 3: SYSTEM ARCHITECTURE
    # ==========================================
    slide3 = add_blank_slide()
    add_header(slide3, "End-to-End System Architecture on AWS", "Decoupling intelligence from execution via AWS Serverless primitives.")

    steps = [
        ("1. Client Application", "Web UI / Rest Client sends prompt via HTTP POST to API Gateway.", ACCENT_BLUE),
        ("2. API Gateway", "Provides REST endpoint, CORS handling, and SSL termination.", ACCENT_BLUE),
        ("3. Bedrock Agent", "Amazon Nova Micro executes reasoning loops & tool discovery.", ACCENT_BEDROCK),
        ("4. MCP Server", "AWS Lambda executes standardized tools/list & tools/call handlers.", ACCENT_MCP),
        ("5. DynamoDB Storage", "Pay-per-request session memory table with automatic 24h TTL.", ACCENT_AWS)
    ]

    for i, (title, desc, color) in enumerate(steps):
        y_pos = Inches(1.8 + (i * 1.0))
        box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(11.7), Inches(0.8))
        box.fill.solid()
        box.fill.fore_color.rgb = BG_CARD
        box.line.color.rgb = color

        tf_s = box.text_frame
        tf_s.margin_left = Inches(0.3)
        tf_s.margin_top = Inches(0.15)

        p1 = tf_s.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = color

        p2 = tf_s.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 4: MCP SPECIFICATION
    # ==========================================
    slide4 = add_blank_slide()
    add_header(slide4, "Model Context Protocol (MCP) Technical Specification", "Understanding JSON-RPC standard payloads for AI Tool Discovery and Execution.")

    # Left: tools/list
    code1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    code1.fill.solid()
    code1.fill.fore_color.rgb = BG_CARD
    code1.line.color.rgb = ACCENT_MCP

    tf_m1 = code1.text_frame
    tf_m1.margin_left = Inches(0.3)
    tf_m1.margin_top = Inches(0.3)
    tf_m1.word_wrap = True

    p = tf_m1.paragraphs[0]
    p.text = "1. Tool Discovery (tools/list)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_MCP

    spec_code1 = """{
  "tools": [
    {
      "name": "query_database",
      "description": "Query session history",
      "input_schema": {
        "type": "object",
        "properties": {
          "session_id": { "type": "string" }
        },
        "required": ["session_id"]
      }
    }
  ]
}"""
    p_code1 = tf_m1.add_paragraph()
    p_code1.text = spec_code1
    p_code1.font.size = Pt(11)
    p_code1.font.color.rgb = RGBColor(56, 189, 248)
    p_code1.space_before = Pt(10)

    # Right: tools/call
    code2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
    code2.fill.solid()
    code2.fill.fore_color.rgb = BG_CARD
    code2.line.color.rgb = ACCENT_AWS

    tf_m2 = code2.text_frame
    tf_m2.margin_left = Inches(0.3)
    tf_m2.margin_top = Inches(0.3)
    tf_m2.word_wrap = True

    p = tf_m2.paragraphs[0]
    p.text = "2. Tool Execution (tools/call)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT_AWS

    spec_code2 = """{
  "action": "tools/call",
  "name": "save_session_data",
  "arguments": {
    "session_id": "session_101",
    "content": "Purchased Prime Upgrade"
  }
}

// Result Output:
{
  "status": "success",
  "message": "Item stored in DynamoDB"
}"""
    p_code2 = tf_m2.add_paragraph()
    p_code2.text = spec_code2
    p_code2.font.size = Pt(11)
    p_code2.font.color.rgb = RGBColor(251, 191, 36)
    p_code2.space_before = Pt(10)

    # ==========================================
    # SLIDE 5: AWS INFRASTRUCTURE TABLE
    # ==========================================
    slide5 = add_blank_slide()
    add_header(slide5, "AWS Infrastructure & Resource Mapping", "Complete CloudFormation serverless stack components.")

    table_shape = slide5.shapes.add_table(6, 4, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    table = table_shape.table
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(3.2)
    table.columns[2].width = Inches(4.3)
    table.columns[3].width = Inches(2.0)

    headers = ["Component", "AWS Resource", "Role & Purpose", "Billing Mode"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG_CARD
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.bold = True
            paragraph.font.color.rgb = ACCENT_AWS

    rows = [
        ("Reasoning Engine", "Amazon Bedrock (Nova Micro)", "Autonomous reasoning and tool synthesis", "Per 1K Tokens"),
        ("Control Loop", "AWS Lambda (Agent Runtime)", "Orchestrates Converse API and tool loops", "Per Millisecond"),
        ("MCP Server", "AWS Lambda (MCP Tool Server)", "Executes database and API logic", "Per Millisecond"),
        ("Session Storage", "Amazon DynamoDB", "Pay-per-request session memory table", "On-Demand"),
        ("Infrastructure", "AWS CloudFormation", "100% IaC stack automated deployment", "Free Tier")
    ]

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i+1, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_DARK if i % 2 == 0 else BG_CARD
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = TEXT_WHITE

    # Save output
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Agentic_AI_AWS_Bedrock_MCP.pptx")
    prs.save(output_path)
    print(f"[PPT_SUCCESS] Presentation saved to: {output_path}")

if __name__ == '__main__':
    create_presentation()
